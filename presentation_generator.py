"""
Генерация PPTX презентаций из JSON-структуры слайдов.
"""
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SLIDE_BG_RGB = (41, 50, 65)  # тёмно-синий
TEXT_COLOR_RGB = (255, 255, 255)


def _set_slide_background(slide, rgb=SLIDE_BG_RGB):
    """Задаёт цветной фон слайда."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _set_word_wrap(tf):
    """Включает автоперенос текста."""
    tf.word_wrap = True


def _add_title_slide(prs, slide_data):
    """Слайд 1: title."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = slide_data.get("title", "Заголовок")
    body.text = slide_data.get("subtitle", "")
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    body.text_frame.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    _set_word_wrap(title.text_frame)
    _set_word_wrap(body.text_frame)


def _add_toc_bullets_slide(prs, slide_data):
    """Слайды 2, 4: toc, bullets."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = slide_data.get("title", "")
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    _set_word_wrap(title.text_frame)
    points = slide_data.get("points", []) or ["—"]
    tf = body.text_frame
    tf.clear()
    for i, p in enumerate(points):
        if i == 0:
            p_para = tf.paragraphs[0]
        else:
            p_para = tf.add_paragraph()
        p_para.text = str(p)
        p_para.font.size = Pt(18)
        p_para.font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    _set_word_wrap(tf)


def _add_text_slide(prs, slide_data):
    """Слайды 3, 7, 8: intro, conclusion, outro."""
    layout = prs.slide_layouts[6]  # пустой
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide)
    left, top, width, height = Inches(0.5), Inches(1), Inches(12.3), Inches(5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = slide_data.get("title", "")
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = slide_data.get("text", "")
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    p2.space_before = Pt(12)
    _set_word_wrap(tf)


def _add_mindmap_slide(prs, slide_data):
    """Слайд 5: mindmap — слева Core, справа Branches в колонку."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide)
    core = slide_data.get("core", "Центр")
    branches = slide_data.get("branches", [])
    # Core: слева по центру (X: 1, Y: 3, W: 3.5, H: 1.5)
    core_box = slide.shapes.add_shape(1, Inches(1), Inches(3), Inches(3.5), Inches(1.5))
    core_box.fill.solid()
    core_box.fill.fore_color.rgb = RGBColor(59, 130, 246)
    core_box.line.fill.background()
    tf = core_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = core
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Branches: справа в колонку (X: 6, Y: 1.5 + i*1.2, W: 6, H: 0.8)
    for i, branch in enumerate(branches):
        y = 1.5 + i * 1.2
        box = slide.shapes.add_shape(1, Inches(6), Inches(y), Inches(6), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(51, 65, 85)
        box.line.fill.background()
        btf = box.text_frame
        btf.word_wrap = True
        btf.paragraphs[0].text = str(branch)[:80]
        btf.paragraphs[0].font.size = Pt(11)
        btf.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    # Заголовок слайда
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = slide_data.get("title", "Mindmap")
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)


def _add_two_columns_slide(prs, slide_data):
    """Слайд 6: two_columns."""
    layout = prs.slide_layouts[3]  # два объекта
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide)
    title = slide.shapes.title
    title.text = slide_data.get("title", "")
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    _set_word_wrap(title.text_frame)
    placeholders = [p for p in slide.placeholders if p.placeholder_format.idx > 1]
    if len(placeholders) >= 2:
        placeholders[0].text = f"{slide_data.get('col1_title', '')}\n\n{slide_data.get('col1_text', '')}"
        placeholders[1].text = f"{slide_data.get('col2_title', '')}\n\n{slide_data.get('col2_text', '')}"
        for ph in placeholders:
            _set_word_wrap(ph.text_frame)
            for p in ph.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
    else:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(4))
        tx.text_frame.text = f"{slide_data.get('col1_title', '')}\n{slide_data.get('col1_text', '')}"
        _set_word_wrap(tx.text_frame)
        for p in tx.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(*TEXT_COLOR_RGB)
        tx2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(4))
        tx2.text_frame.text = f"{slide_data.get('col2_title', '')}\n{slide_data.get('col2_text', '')}"
        _set_word_wrap(tx2.text_frame)
        for p in tx2.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(*TEXT_COLOR_RGB)


def build_pptx(slides_data: list[dict], output_path: Path) -> None:
    """Создаёт PPTX из списка слайдов. Формат 16:9."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for s in slides_data:
        stype = s.get("type", "intro")
        if stype == "title":
            _add_title_slide(prs, s)
        elif stype in ("toc", "bullets"):
            _add_toc_bullets_slide(prs, s)
        elif stype == "mindmap":
            _add_mindmap_slide(prs, s)
        elif stype == "two_columns":
            _add_two_columns_slide(prs, s)
        else:
            _add_text_slide(prs, s)
    prs.save(output_path)


def generate_presentation_filename() -> str:
    """Уникальное имя для презентации (без расширения)."""
    return f"pres_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
